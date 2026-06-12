"""
WishKeeper Investor Pitch Deck — python-pptx generator
Berry & Cream palette: Berry #6D2E46, Dusty Rose #A26769, Cream #ECE2D0
Motif: rounded icon circles on every content slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ───────────────────────────────────────────────────────────────────
BERRY     = RGBColor(0x6D, 0x2E, 0x46)
ROSE      = RGBColor(0xA2, 0x67, 0x69)
CREAM     = RGBColor(0xEC, 0xE2, 0xD0)

def rgb(r,g,b): return RGBColor(r,g,b)

TEAL_C    = rgb(0x1D, 0x9E, 0x75)
LIGHT_ROSE= RGBColor(0xF5, 0xED, 0xE8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x0A, 0x12)
MID_GRAY  = RGBColor(0x88, 0x70, 0x78)
GOLD      = RGBColor(0xC8, 0x9B, 0x5A)

def rgb(r,g,b): return RGBColor(r,g,b)

# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None, radius=None):
    from pptx.util import Pt as Pt2
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width: shape.line.width = Pt2(line_width)
    else:
        shape.line.fill.background()
    if radius:
        # Apply rounded corners via XML
        sp = shape._element
        sp_pr = sp.find(qn('p:spPr'))
        prstGeom = sp_pr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.attrib['prst'] = 'roundRect'
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                from lxml import etree
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            from lxml import etree
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.attrib['name'] = 'adj'
            gd.attrib['fmla'] = f'val {int(radius * 100000)}'
    return shape

def txt(slide, text, x, y, w, h, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False, font='Calibri'):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb

def circle(slide, cx, cy, r, fill_color):
    """Draw a filled circle centered at (cx,cy) with radius r (inches)"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        9,  # oval
        Inches(cx - r), Inches(cy - r), Inches(r*2), Inches(r*2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def divider(slide, x, y, w, color=ROSE, thickness=0.03):
    box(slide, x, y, w, thickness, fill_color=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, BERRY)

# Soft cream arc shape (decorative circle)
shape = s.shapes.add_shape(9, Inches(8.5), Inches(-1.5), Inches(8), Inches(8))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x82, 0x3A, 0x58)
shape.line.fill.background()

shape2 = s.shapes.add_shape(9, Inches(10), Inches(3.5), Inches(5), Inches(5))
shape2.fill.solid(); shape2.fill.fore_color.rgb = RGBColor(0x78, 0x33, 0x50)
shape2.line.fill.background()

# Gift icon circle
c = circle(s, 2.5, 2.4, 0.6, ROSE)
txt(s, "🎁", 1.95, 1.9, 1.1, 1.1, size=28, align=PP_ALIGN.CENTER)

txt(s, "WishKeeper", 0.8, 3.3, 7, 1.2, size=52, bold=True, color=WHITE, font='Cambria')
txt(s, "The celebration platform your family has been waiting for.",
    0.8, 4.55, 8, 0.7, size=20, color=CREAM, italic=True, font='Calibri')

divider(s, 0.8, 5.4, 3.5, color=GOLD)

txt(s, "Seed Round  ·  2026", 0.8, 5.6, 5, 0.5, size=14, color=ROSE, font='Calibri')
txt(s, "Confidential — not for distribution", 9.5, 7.0, 3.5, 0.35, size=10, color=ROSE, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "The problem", 0.6, 0.25, 8, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "Families want to celebrate — but don't know how", 0.6, 1.0, 10, 0.5, size=16, color=CREAM, font='Calibri')

pain_points = [
    ("😞", "\"What do you want for your birthday?\" — \"Oh, nothing\""),
    ("🎁", "Duplicate gifts, wrong sizes, generic choices"),
    ("📋", "No single place to capture how someone wants to be celebrated"),
    ("📱", "Existing apps (Giftster, Elfster) focus on products, not experiences or milestones"),
    ("💔", "Milestone birthdays — 40th, 50th, 60th — go under-celebrated"),
]
for i, (emoji, text) in enumerate(pain_points):
    y = 2.0 + i * 0.9
    circle(s, 1.1, y + 0.22, 0.28, LIGHT_ROSE)
    txt(s, emoji, 0.86, y, 0.5, 0.5, size=16, align=PP_ALIGN.CENTER)
    txt(s, text, 1.65, y, 10.8, 0.55, size=16, color=DARK, font='Calibri')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT_ROSE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "The solution", 0.6, 0.25, 8, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "WishKeeper — \"How I Want to Be Celebrated\"", 0.6, 0.95, 11, 0.55, size=16, color=CREAM, italic=True, font='Calibri')

features = [
    ("🎂", "Celebration Profiles", "Birthday vision boards, anniversary dreams, Eid traditions"),
    ("🎁", "Wish Lists", "Gifts, experiences, restaurants, jewelry — with prices and links"),
    ("👨‍👩‍👧", "Family Access Controls", "Husband sees everything; kids see gifts only — you decide"),
    ("🔒", "Secret Gift Tracker", "Family coordinates gifts without spoiling surprises"),
    ("💌", "Memory Vault", "Photos, letters, recipes, and time-locked letters to loved ones"),
    ("✨", "AI Assistant", "Personalized gift ideas, party themes, and celebration writing"),
]

for i, (emoji, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.55 + col * 6.4
    y = 2.0 + row * 1.55
    box(s, x, y, 5.9, 1.35, fill_color=WHITE, radius=0.08)
    circle(s, x + 0.5, y + 0.45, 0.32, BERRY)
    txt(s, emoji, x + 0.22, y + 0.18, 0.6, 0.6, size=16, align=PP_ALIGN.CENTER)
    txt(s, title, x + 1.0, y + 0.1, 4.6, 0.45, size=14, bold=True, color=BERRY, font='Calibri')
    txt(s, desc, x + 1.0, y + 0.55, 4.6, 0.65, size=12, color=MID_GRAY, font='Calibri')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Market Opportunity
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "Market opportunity", 0.6, 0.25, 10, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "A large, underserved market hiding in plain sight", 0.6, 1.0, 10, 0.5, size=16, color=CREAM, font='Calibri')

# Three market circles
markets = [
    ("$8.3B", "Online gift\nmarket (US)", BERRY),
    ("$2.1B", "Event planning\napps (global)", ROSE),
    ("$680M", "Digital legacy\n& memory apps", GOLD),
]
for i, (val, label, color) in enumerate(markets):
    x = 1.2 + i * 3.8
    circle(s, x + 1.0, 3.6, 1.15, color)
    txt(s, val, x - 0.15, 3.0, 2.3, 0.8, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Cambria')
    txt(s, label, x - 0.15, 3.85, 2.3, 0.7, size=12, color=WHITE, align=PP_ALIGN.CENTER, font='Calibri')

txt(s, "Target: 48M US women aged 35–65 who manage family celebrations",
    0.6, 5.3, 12, 0.55, size=15, color=DARK, font='Calibri', align=PP_ALIGN.CENTER)
txt(s, "Primary persona: mothers, wives, and matriarchs in multigenerational families",
    0.6, 5.85, 12, 0.5, size=14, color=MID_GRAY, italic=True, font='Calibri', align=PP_ALIGN.CENTER)

divider(s, 0.6, 6.5, 12.1, color=CREAM)
txt(s, "Serviceable Addressable Market: $340M at $7.99/mo across 3.5M paying users",
    0.6, 6.65, 12.1, 0.5, size=13, color=MID_GRAY, font='Calibri', align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Competitive Landscape
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT_ROSE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "Competitive landscape", 0.6, 0.25, 10, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "No one owns \"How I Want to Be Celebrated\"", 0.6, 1.0, 10, 0.5, size=16, color=CREAM, italic=True, font='Calibri')

headers = ["Product focus", "Experiences", "Event planning", "Family traditions", "Memory vault"]
cols = [("Giftster", ROSE), ("Elfster", ROSE), ("Pinterest", ROSE), ("WishKeeper", BERRY)]

# Header row
for j, h in enumerate(headers):
    x = 2.6 + j * 2.1
    txt(s, h, x, 2.0, 2.0, 0.4, size=12, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER, font='Calibri')

check_data = {
    "Giftster":    ["✓", "✗", "✗", "✗", "✗"],
    "Elfster":     ["✓", "✗", "✗", "✗", "✗"],
    "Pinterest":   ["✗", "✓", "✓", "✗", "✗"],
    "WishKeeper":  ["✓", "✓", "✓", "✓", "✓"],
}

for i, (name, color) in enumerate(cols):
    y = 2.6 + i * 0.9
    is_wish = (name == "WishKeeper")
    row_bg = BERRY if is_wish else WHITE
    row_text = WHITE if is_wish else DARK
    box(s, 0.5, y - 0.1, 12.3, 0.75, fill_color=row_bg, radius=0.06)
    txt(s, name, 0.7, y, 1.8, 0.55, size=14, bold=is_wish, color=row_text, font='Calibri')
    checks = check_data[name]
    for j, ch in enumerate(checks):
        x = 2.6 + j * 2.1
        chk_color = GOLD if (ch == "✓" and is_wish) else (TEAL_C if ch == "✓" else LIGHT_ROSE)
        if not is_wish:
            chk_color = rgb(0x1D,0x9E,0x75) if ch == "✓" else rgb(0xFF, 0xCC, 0xCC)
        txt(s, ch, x + 0.5, y, 1.0, 0.55, size=16, bold=is_wish, color=GOLD if is_wish and ch=="✓" else (rgb(0x1D,0x9E,0x75) if ch=="✓" else rgb(0xBB,0x44,0x44)), align=PP_ALIGN.CENTER, font='Calibri')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Business Model
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "Business model", 0.6, 0.25, 8, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "Simple SaaS — three tiers, low churn by design", 0.6, 1.0, 10, 0.5, size=16, color=CREAM, font='Calibri')

plans = [
    ("Free", "$0/mo", CREAM, DARK, ["1 profile", "Basic wish lists", "Up to 3 family members", "Core reminders"]),
    ("Premium", "$7.99/mo", ROSE, WHITE, ["Unlimited wishes & events", "AI celebration assistant", "Memory vault", "Secret gift tracker", "Priority reminders"]),
    ("Family Plan", "$14.99/mo", BERRY, WHITE, ["Everything in Premium", "Up to 10 family members", "Shared celebration calendar", "Collaborative planning", "Family admin dashboard"]),
]

for i, (name, price, bg_color, fg_color, features_list) in enumerate(plans):
    x = 0.5 + i * 4.25
    is_featured = (name == "Premium")
    box(s, x, 2.0, 3.9, 4.9, fill_color=bg_color, radius=0.1)
    if is_featured:
        badge_box = s.shapes.add_shape(1, Inches(x+1.2), Inches(1.72), Inches(1.5), Inches(0.35))
        badge_box.fill.solid(); badge_box.fill.fore_color.rgb = GOLD
        badge_box.line.fill.background()
        txt(s, "Most popular", x+1.05, 1.75, 1.8, 0.3, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER, font='Calibri')
    txt(s, name, x+0.2, 2.15, 3.5, 0.55, size=20, bold=True, color=fg_color, font='Cambria')
    txt(s, price, x+0.2, 2.75, 3.5, 0.6, size=28, bold=True, color=fg_color if name!="Free" else BERRY, font='Cambria')
    for j, feat in enumerate(features_list):
        txt(s, "✓  " + feat, x+0.25, 3.5 + j*0.5, 3.4, 0.45, size=12, color=fg_color, font='Calibri')

txt(s, "Projected blended ARPU: $6.20/mo  ·  Target Year-3 paid users: 280,000  ·  ARR target: $20.8M",
    0.6, 7.1, 12.1, 0.4, size=12, color=MID_GRAY, font='Calibri', align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Growth Strategy
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT_ROSE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "Growth strategy", 0.6, 0.25, 8, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "Viral by design — every invite is a new user", 0.6, 1.0, 10, 0.5, size=16, color=CREAM, font='Calibri')

channels = [
    ("1", "Viral family invites", "Every profile owner invites 4–6 family members on average. Each invite email drives a signup.\nK-factor: 0.6–0.8 projected."),
    ("2", "Muslim & Arab-American communities", "Eid gift-giving is a $4B+ occasion in the US. No app serves this community's celebration needs well. Partner with Muslim influencers and community organizations."),
    ("3", "Milestone birthday events", "Women turning 40, 50, and 60 actively seek planning tools. Target through Facebook groups, gift guides, and wedding/event planners."),
    ("4", "LegacyCare integration", "Companion product synergy — cross-promote to LegacyCare's existing user base (continuity of celebration legacy)."),
]

for i, (num, title, desc) in enumerate(channels):
    y = 2.0 + i * 1.3
    box(s, 0.5, y, 12.3, 1.15, fill_color=WHITE, radius=0.08)
    circle(s, 1.0, y + 0.42, 0.32, BERRY)
    txt(s, num, 0.75, y + 0.15, 0.5, 0.55, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Cambria')
    txt(s, title, 1.6, y + 0.05, 3.5, 0.45, size=14, bold=True, color=BERRY, font='Calibri')
    txt(s, desc, 1.6, y + 0.52, 10.8, 0.55, size=11, color=MID_GRAY, font='Calibri')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Financial Projections
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "Financial projections", 0.6, 0.25, 8, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "Conservative base case — strong unit economics from day one", 0.6, 1.0, 11, 0.5, size=16, color=CREAM, font='Calibri')

# Table
headers_t = ["", "Year 1", "Year 2", "Year 3"]
rows_t = [
    ["Registered users",    "22,000",   "95,000",   "380,000"],
    ["Paid subscribers",    "3,200",    "28,000",   "112,000"],
    ["ARPU (monthly)",      "$5.40",    "$6.10",    "$6.80"],
    ["ARR",                 "$207K",    "$2.05M",   "$9.1M"],
    ["Gross margin",        "82%",      "84%",      "86%"],
    ["CAC",                 "$18",      "$14",      "$11"],
    ["LTV (36mo)",          "$194",     "$220",     "$245"],
    ["LTV:CAC",             "10.8x",    "15.7x",    "22.3x"],
]

for j, h in enumerate(headers_t):
    x = 0.5 + j * 3.0
    box(s, x, 2.1, 2.9, 0.5, fill_color=BERRY)
    txt(s, h, x+0.1, 2.15, 2.7, 0.4, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Calibri')

for i, row in enumerate(rows_t):
    row_bg = LIGHT_ROSE if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        x = 0.5 + j * 3.0
        box(s, x, 2.6 + i*0.5, 2.9, 0.5, fill_color=row_bg)
        align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        weight = j == 0
        txt(s, cell, x+0.1, 2.65+i*0.5, 2.7, 0.4, size=12, bold=weight, color=DARK, align=align, font='Calibri')

txt(s, "* Projections assume $1.5M seed, 18-month runway to Series A", 0.6, 7.1, 12, 0.4, size=11, color=MID_GRAY, italic=True, font='Calibri')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Traction & Roadmap
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT_ROSE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "Traction & roadmap", 0.6, 0.25, 8, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "Built on validated demand, clear path to scale", 0.6, 1.0, 10, 0.5, size=16, color=CREAM, font='Calibri')

# Left: Traction
box(s, 0.5, 2.0, 5.8, 5.2, fill_color=WHITE, radius=0.1)
txt(s, "Current traction", 0.9, 2.2, 5.0, 0.5, size=16, bold=True, color=BERRY, font='Cambria')
traction = [
    ("✓", "MVP built with full React Native codebase"),
    ("✓", "Firebase backend + security rules complete"),
    ("✓", "AI assistant integrated (Claude API)"),
    ("✓", "Beta waitlist: 340 sign-ups (organic)"),
    ("✓", "User research: 28 interviews with target demographic"),
    ("✓", "NPS from concept tests: 71"),
    ("✓", "LegacyCare founder partnership confirmed"),
]
for i, (icon, line) in enumerate(traction):
    txt(s, icon + "  " + line, 0.85, 2.9 + i*0.47, 5.2, 0.42, size=12, color=DARK, font='Calibri')

# Right: Roadmap
box(s, 6.8, 2.0, 6.0, 5.2, fill_color=WHITE, radius=0.1)
txt(s, "12-month roadmap", 7.2, 2.2, 5.2, 0.5, size=16, bold=True, color=BERRY, font='Cambria')
milestones = [
    ("Q3 2026", "Beta launch (iOS). 500 invited users. Gift tracker, wishes, events."),
    ("Q4 2026", "Public launch iOS + Android. Paid tiers live. Target 3,000 paid."),
    ("Q1 2027", "AI personalization v2. Pinterest-style boards. Group gifting."),
    ("Q2 2027", "Web family dashboard. Etsy/Amazon integrations. Series A prep."),
]
for i, (q, desc) in enumerate(milestones):
    y = 2.9 + i * 1.1
    circle(s, 7.5, y + 0.22, 0.22, BERRY)
    txt(s, q, 7.18, y, 1.0, 0.42, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Calibri')
    txt(s, desc, 7.95, y, 4.5, 0.85, size=12, color=DARK, font='Calibri')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — The Ask
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, BERRY)

shape = s.shapes.add_shape(9, Inches(8), Inches(-2), Inches(9), Inches(9))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x82, 0x3A, 0x58)
shape.line.fill.background()

txt(s, "The ask", 0.8, 0.9, 9, 1.0, size=44, bold=True, color=WHITE, font='Cambria')
txt(s, "Seed Round  ·  $1.5M", 0.8, 1.9, 9, 0.8, size=32, bold=True, color=GOLD, font='Cambria')

divider(s, 0.8, 2.85, 5, color=ROSE)

uses = [
    ("40%", "$600K", "Product & Engineering (12-month build-out)"),
    ("30%", "$450K", "Growth & Community (beta → public launch)"),
    ("20%", "$300K", "Team (2 engineers + 1 designer)"),
    ("10%", "$150K", "Operations, legal, infrastructure"),
]
for i, (pct, amt, label) in enumerate(uses):
    y = 3.2 + i * 0.85
    circle(s, 1.3, y + 0.22, 0.28, ROSE)
    txt(s, pct, 1.05, y + 0.02, 0.5, 0.42, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Calibri')
    txt(s, amt, 1.85, y, 1.2, 0.45, size=16, bold=True, color=GOLD, font='Cambria')
    txt(s, label, 3.1, y + 0.05, 8, 0.42, size=14, color=CREAM, font='Calibri')

txt(s, "18-month runway  ·  Target: 28K paid subscribers  ·  Series A ready at $2M ARR",
    0.8, 7.0, 11, 0.4, size=13, color=ROSE, font='Calibri')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Team
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.6, fill_color=BERRY)

txt(s, "The team", 0.6, 0.25, 8, 0.8, size=36, bold=True, color=WHITE, font='Cambria')
txt(s, "Founders who've built in this space before", 0.6, 1.0, 10, 0.5, size=16, color=CREAM, font='Calibri')

team = [
    ("Founder / CEO", "LegacyCare founder. Deep expertise in family continuity apps. 0→1 product builder. Target demographic: lived experience."),
    ("CTO (Hiring)", "React Native + Firebase specialist. 5+ yrs building consumer mobile apps."),
    ("Head of Growth (Hiring)", "Community-led growth expert. Muslim/Arab-American market experience a plus."),
    ("Advisor", "Exits in SaaS and consumer subscription. Angel network access."),
]
for i, (role, bio) in enumerate(team):
    col = i % 2; row = i // 2
    x = 0.5 + col * 6.4; y = 2.0 + row * 2.3
    box(s, x, y, 6.0, 2.1, fill_color=LIGHT_ROSE, radius=0.1)
    circle(s, x + 0.55, y + 0.55, 0.42, BERRY)
    txt(s, "👤", x + 0.26, y + 0.25, 0.6, 0.6, size=20, align=PP_ALIGN.CENTER)
    txt(s, role, x + 1.2, y + 0.15, 4.5, 0.5, size=14, bold=True, color=BERRY, font='Cambria')
    txt(s, bio, x + 1.2, y + 0.65, 4.5, 1.2, size=12, color=MID_GRAY, font='Calibri')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Close
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, BERRY)

shape = s.shapes.add_shape(9, Inches(-2), Inches(3), Inches(8), Inches(8))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x82, 0x3A, 0x58)
shape.line.fill.background()

txt(s, "Every woman deserves to be\ncelebrated exactly as she dreams.", 1.5, 1.4, 10.3, 2.0, size=32, bold=True, color=WHITE, font='Cambria', align=PP_ALIGN.CENTER)

divider(s, 4.5, 3.6, 4.3, color=GOLD)

txt(s, "WishKeeper is how that happens.", 1.5, 3.9, 10.3, 0.7, size=22, italic=True, color=CREAM, font='Calibri', align=PP_ALIGN.CENTER)

txt(s, "hello@wishkeeper.app  ·  wishkeeper.app/investors", 1.5, 5.5, 10.3, 0.55, size=15, color=ROSE, align=PP_ALIGN.CENTER, font='Calibri')
txt(s, "Seed Round — $1.5M  ·  Deck version June 2026", 1.5, 6.15, 10.3, 0.45, size=12, color=ROSE, align=PP_ALIGN.CENTER, font='Calibri')
txt(s, "Confidential — not for distribution", 0.5, 7.1, 12.3, 0.35, size=10, color=ROSE, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "/mnt/user-data/outputs/WishKeeper_Investor_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
